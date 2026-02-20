#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT Processor Skill - PowerPoint演示文稿处理工具
支持从Markdown/JSON创建PPT、模板应用、图表插入和图片处理

Author: godlike-kimi
Version: 1.0.0
License: MIT
"""

import argparse
import json
import os
import re
import sys
from typing import Dict, List, Optional, Tuple, Any
from io import BytesIO

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from pptx.oxml import parse_xml
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    Image = None

try:
    import matplotlib.pyplot as plt
    import matplotlib
    matplotlib.use('Agg')
except ImportError:
    plt = None


class PPTXSkill:
    """PowerPoint演示文稿处理类"""
    
    # 预设主题配置
    THEMES = {
        "default": {
            "bg_color": (255, 255, 255),
            "title_color": (0, 0, 0),
            "text_color": (51, 51, 51),
            "accent_color": (0, 112, 192)
        },
        "dark": {
            "bg_color": (40, 40, 40),
            "title_color": (255, 255, 255),
            "text_color": (220, 220, 220),
            "accent_color": (0, 150, 255)
        },
        "light": {
            "bg_color": (245, 245, 245),
            "title_color": (50, 50, 50),
            "text_color": (80, 80, 80),
            "accent_color": (70, 130, 180)
        },
        "blue": {
            "bg_color": (240, 248, 255),
            "title_color": (0, 51, 102),
            "text_color": (51, 51, 102),
            "accent_color": (0, 102, 204)
        },
        "green": {
            "bg_color": (240, 255, 240),
            "title_color": (0, 80, 0),
            "text_color": (34, 85, 51),
            "accent_color": (34, 139, 34)
        }
    }
    
    # 幻灯片布局映射
    LAYOUTS = {
        "title": 0,
        "title_and_content": 1,
        "section_header": 2,
        "two_content": 3,
        "comparison": 4,
        "title_only": 5,
        "blank": 6,
        "content_with_caption": 7,
        "picture_with_caption": 8
    }
    
    def __init__(self, template_path: Optional[str] = None):
        """
        初始化PPT处理器
        
        Args:
            template_path: 模板文件路径，None则创建空白演示文稿
        """
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        self.current_theme = "default"
        self._apply_default_size()
    
    def _apply_default_size(self):
        """设置默认幻灯片尺寸为16:9"""
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
    
    def set_theme(self, theme_name: str):
        """
        设置主题
        
        Args:
            theme_name: 主题名称
        """
        if theme_name in self.THEMES:
            self.current_theme = theme_name
    
    def _get_theme_color(self, color_type: str) -> RGBColor:
        """获取主题颜色"""
        theme = self.THEMES.get(self.current_theme, self.THEMES["default"])
        color_tuple = theme.get(color_type, (0, 0, 0))
        return RGBColor(*color_tuple)
    
    def add_slide(self, title: str = "", content: List[str] = None, 
                  layout: str = "title_and_content", subtitle: str = "") -> Any:
        """
        添加幻灯片
        
        Args:
            title: 幻灯片标题
            content: 内容列表
            layout: 布局类型
            subtitle: 副标题
            
        Returns:
            幻灯片对象
        """
        layout_idx = self.LAYOUTS.get(layout, 1)
        slide_layout = self.prs.slide_layouts[layout_idx]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 设置标题
        if title and slide.shapes.title:
            slide.shapes.title.text = title
            title_frame = slide.shapes.title.text_frame
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = self._get_theme_color("title_color")
        
        # 设置内容
        if content and len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            
            for i, item in enumerate(content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(18)
                p.font.color.rgb = self._get_theme_color("text_color")
                p.space_after = Pt(12)
        
        return slide
    
    def add_title_slide(self, title: str, subtitle: str = ""):
        """添加标题幻灯片"""
        slide_layout = self.prs.slide_layouts[0]  # 标题幻灯片布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            title_frame = slide.shapes.title.text_frame
            title_frame.paragraphs[0].font.size = Pt(44)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = self._get_theme_color("title_color")
        
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
            subtitle_frame = slide.placeholders[1].text_frame
            subtitle_frame.paragraphs[0].font.size = Pt(24)
            subtitle_frame.paragraphs[0].font.color.rgb = self._get_theme_color("text_color")
    
    def add_image_slide(self, title: str, image_path: str, 
                        left: float = 1, top: float = 1.5, 
                        width: Optional[float] = None, height: Optional[float] = None):
        """
        添加带图片的幻灯片
        
        Args:
            title: 幻灯片标题
            image_path: 图片路径
            left, top: 图片位置（英寸）
            width, height: 图片尺寸（英寸），None则自动调整
        """
        slide = self.add_slide(title=title, layout="title_only")
        
        if not os.path.exists(image_path):
            print(f"Warning: Image not found: {image_path}")
            return
        
        # 自动调整图片大小
        if width is None or height is None:
            try:
                with Image.open(image_path) as img:
                    img_width, img_height = img.size
                    aspect_ratio = img_width / img_height
                    
                    if width is None and height is None:
                        width = 8
                        height = width / aspect_ratio
                    elif width is None:
                        width = height * aspect_ratio
                    else:
                        height = width / aspect_ratio
            except Exception as e:
                print(f"Warning: Could not process image: {e}")
                width, height = 6, 4
        
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), 
                                  Inches(width), Inches(height))
    
    def add_chart(self, slide_index: int, chart_type: str, 
                  data: Dict[str, List], title: str = ""):
        """
        添加图表到幻灯片
        
        Args:
            slide_index: 幻灯片索引（从1开始）
            chart_type: 图表类型 (bar, line, pie)
            data: 图表数据 {"labels": [...], "values": [...]}
            title: 图表标题
        """
        if slide_index < 1 or slide_index > len(self.prs.slides):
            print(f"Error: Invalid slide index {slide_index}")
            return
        
        if plt is None:
            print("Warning: matplotlib not installed, skipping chart")
            return
        
        slide = self.prs.slides[slide_index - 1]
        labels = data.get("labels", [])
        values = data.get("values", [])
        
        # 创建matplotlib图表
        fig, ax = plt.subplots(figsize=(6, 4))
        
        if chart_type == "bar":
            ax.bar(labels, values, color=self._get_theme_color("accent_color"))
        elif chart_type == "line":
            ax.plot(labels, values, marker='o', linewidth=2, 
                   color=self._get_theme_color("accent_color"))
        elif chart_type == "pie":
            ax.pie(values, labels=labels, autopct='%1.1f%%')
        
        if title:
            ax.set_title(title)
        
        # 保存到内存
        img_stream = BytesIO()
        plt.savefig(img_stream, format='png', dpi=150, bbox_inches='tight')
        img_stream.seek(0)
        plt.close(fig)
        
        # 添加到幻灯片
        slide.shapes.add_picture(img_stream, Inches(1), Inches(2), Inches(8))
    
    def parse_markdown(self, markdown_text: str) -> List[Dict]:
        """
        解析Markdown文本为幻灯片数据
        
        Args:
            markdown_text: Markdown格式文本
            
        Returns:
            幻灯片数据列表
        """
        slides_data = []
        # 按 --- 分割幻灯片
        slides_raw = re.split(r'\n---\s*\n', markdown_text)
        
        for slide_raw in slides_raw:
            slide_raw = slide_raw.strip()
            if not slide_raw:
                continue
            
            slide_data = {"title": "", "content": [], "images": []}
            lines = slide_raw.split('\n')
            current_section = "content"
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # 解析标题
                if line.startswith('# '):
                    slide_data["title"] = line[2:].strip()
                elif line.startswith('## '):
                    slide_data["title"] = line[3:].strip()
                # 解析列表项
                elif line.startswith('- ') or line.startswith('* '):
                    slide_data["content"].append(line[2:].strip())
                elif re.match(r'^\d+\.', line):
                    content = re.sub(r'^\d+\.\s*', '', line)
                    slide_data["content"].append(content)
                # 解析图片
                elif line.startswith('!['):
                    img_match = re.match(r'!\[([^\]]*)\]\(([^)]+)\)', line)
                    if img_match:
                        slide_data["images"].append(img_match.group(2))
                # 其他内容
                else:
                    slide_data["content"].append(line)
            
            slides_data.append(slide_data)
        
        return slides_data
    
    def create_from_markdown(self, markdown_text: str, output_path: str, 
                             template: str = "default", title_slide: bool = True):
        """
        从Markdown创建PPT
        
        Args:
            markdown_text: Markdown文本或文件路径
            output_path: 输出文件路径
            template: 主题名称
            title_slide: 是否添加标题页
        """
        self.set_theme(template)
        
        # 如果是文件路径，读取内容
        if os.path.exists(markdown_text):
            with open(markdown_text, 'r', encoding='utf-8') as f:
                markdown_text = f.read()
        
        slides_data = self.parse_markdown(markdown_text)
        
        # 添加标题页
        if title_slide and slides_data:
            first_slide = slides_data[0]
            self.add_title_slide(
                first_slide.get("title", "Presentation"),
                ""
            )
            slides_data = slides_data[1:]
        
        # 添加内容幻灯片
        for slide_data in slides_data:
            self.add_slide(
                title=slide_data.get("title", ""),
                content=slide_data.get("content", [])
            )
        
        self.save(output_path)
    
    def create_from_json(self, json_data: str, output_path: str, template: str = "default"):
        """
        从JSON创建PPT
        
        Args:
            json_data: JSON字符串或文件路径
            output_path: 输出文件路径
            template: 主题名称
        """
        self.set_theme(template)
        
        # 如果是文件路径，读取内容
        if os.path.exists(json_data):
            with open(json_data, 'r', encoding='utf-8') as f:
                data = json.load(f)
        else:
            data = json.loads(json_data)
        
        # 添加标题页
        title = data.get("title", "Presentation")
        self.add_title_slide(title, data.get("subtitle", ""))
        
        # 添加幻灯片
        for slide_data in data.get("slides", []):
            slide = self.add_slide(
                title=slide_data.get("title", ""),
                content=slide_data.get("content", []),
                layout=slide_data.get("layout", "title_and_content")
            )
            
            # 添加图表
            if "chart" in slide_data:
                chart = slide_data["chart"]
                self.add_chart(
                    len(self.prs.slides),
                    chart.get("type", "bar"),
                    chart.get("data", {}),
                    chart.get("title", "")
                )
        
        self.save(output_path)
    
    def merge_presentations(self, pptx_files: List[str], output_path: str):
        """
        合并多个PPT文件
        
        Args:
            pptx_files: PPT文件路径列表
            output_path: 输出文件路径
        """
        merged_prs = Presentation()
        merged_prs.slide_width = Inches(13.333)
        merged_prs.slide_height = Inches(7.5)
        
        for pptx_file in pptx_files:
            if not os.path.exists(pptx_file):
                print(f"Warning: File not found: {pptx_file}")
                continue
            
            source_prs = Presentation(pptx_file)
            for slide in source_prs.slides:
                # 复制幻灯片到新演示文稿
                blank_layout = merged_prs.slide_layouts[6]  # 空白布局
                new_slide = merged_prs.slides.add_slide(blank_layout)
                
                # 复制形状
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        
                        new_shape = new_slide.shapes.add_textbox(left, top, width, height)
                        new_frame = new_shape.text_frame
                        new_frame.text = shape.text_frame.text
        
        merged_prs.save(output_path)
        print(f"Merged presentation saved to: {output_path}")
    
    def save(self, output_path: str):
        """
        保存演示文稿
        
        Args:
            output_path: 输出文件路径
        """
        # 确保目录存在
        os.makedirs(os.path.dirname(os.path.abspath(output_path)) if os.path.dirname(output_path) else '.', exist_ok=True)
        self.prs.save(output_path)
        print(f"Presentation saved to: {output_path}")


def main():
    """命令行入口"""
    parser = argparse.ArgumentParser(description='PPT Processor - PowerPoint演示文稿处理工具')
    parser.add_argument('--action', '-a', required=True,
                       choices=['create', 'edit', 'merge', 'split', 'template'],
                       help='操作类型')
    parser.add_argument('--input', '-i', help='输入文件路径或内容')
    parser.add_argument('--output', '-o', default='output.pptx', help='输出文件路径')
    parser.add_argument('--template', '-t', default='default', help='模板名称')
    parser.add_argument('--theme', default='default', help='主题名称')
    parser.add_argument('--format', '-f', default='markdown', 
                       choices=['markdown', 'json'], help='输入格式')
    
    args = parser.parse_args()
    
    skill = PPTXSkill()
    skill.set_theme(args.theme)
    
    if args.action == 'create':
        if args.format == 'markdown':
            skill.create_from_markdown(args.input or '', args.output, args.template)
        else:
            skill.create_from_json(args.input or '{}', args.output, args.template)
    
    elif args.action == 'merge':
        if not args.input:
            print("Error: --input required for merge action (comma-separated file list)")
            sys.exit(1)
        files = [f.strip() for f in args.input.split(',')]
        skill.merge_presentations(files, args.output)
    
    elif args.action == 'template':
        # 列出可用模板
        print("Available templates:")
        for name in skill.THEMES.keys():
            print(f"  - {name}")
    
    else:
        print(f"Action '{args.action}' not fully implemented yet")


if __name__ == '__main__':
    main()
