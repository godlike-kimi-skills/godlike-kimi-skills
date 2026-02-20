#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT Processor Skill - 基础测试套件
"""

import os
import sys
import json
import tempfile
import unittest
from pathlib import Path

# 添加父目录到路径
sys.path.insert(0, str(Path(__file__).parent.parent))

from main import PPTXSkill

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class TestPPTXSkill(unittest.TestCase):
    """PPT处理器测试类"""
    
    @classmethod
    def setUpClass(cls):
        """测试类初始化"""
        cls.test_dir = tempfile.mkdtemp()
        cls.sample_markdown = """# 测试演示文稿

这是第一页的内容

- 要点1
- 要点2

---

## 第二页

更多内容在这里

---

## 第三页

1. 第一项
2. 第二项
3. 第三项
"""
        cls.sample_json = json.dumps({
            "title": "JSON测试",
            "subtitle": "副标题",
            "slides": [
                {
                    "title": "第一页",
                    "content": ["内容1", "内容2"],
                    "layout": "title_and_content"
                },
                {
                    "title": "图表页",
                    "content": ["数据展示"],
                    "chart": {
                        "type": "bar",
                        "data": {"labels": ["A", "B", "C"], "values": [10, 20, 30]},
                        "title": "示例图表"
                    }
                }
            ]
        })
    
    @classmethod
    def tearDownClass(cls):
        """测试类清理"""
        import shutil
        shutil.rmtree(cls.test_dir, ignore_errors=True)
    
    def setUp(self):
        """每个测试前初始化"""
        if not PPTX_AVAILABLE:
            self.skipTest("python-pptx not installed")
        self.skill = PPTXSkill()
    
    def test_init(self):
        """测试初始化"""
        self.assertIsNotNone(self.skill.prs)
        self.assertEqual(self.skill.current_theme, "default")
    
    def test_set_theme(self):
        """测试主题设置"""
        self.skill.set_theme("dark")
        self.assertEqual(self.skill.current_theme, "dark")
        
        self.skill.set_theme("blue")
        self.assertEqual(self.skill.current_theme, "blue")
        
        # 无效主题应该保持原样
        self.skill.set_theme("invalid_theme")
        self.assertEqual(self.skill.current_theme, "blue")
    
    def test_add_slide(self):
        """测试添加幻灯片"""
        initial_count = len(self.skill.prs.slides)
        
        slide = self.skill.add_slide(
            title="测试标题",
            content=["项目1", "项目2", "项目3"]
        )
        
        self.assertEqual(len(self.skill.prs.slides), initial_count + 1)
        self.assertIsNotNone(slide)
    
    def test_add_title_slide(self):
        """测试添加标题幻灯片"""
        initial_count = len(self.skill.prs.slides)
        
        self.skill.add_title_slide("主标题", "副标题")
        
        self.assertEqual(len(self.skill.prs.slides), initial_count + 1)
    
    def test_parse_markdown(self):
        """测试Markdown解析"""
        slides_data = self.skill.parse_markdown(self.sample_markdown)
        
        self.assertEqual(len(slides_data), 3)
        self.assertEqual(slides_data[0]["title"], "测试演示文稿")
        self.assertIn("要点1", slides_data[0]["content"])
        self.assertEqual(slides_data[1]["title"], "第二页")
    
    def test_create_from_markdown_string(self):
        """测试从Markdown字符串创建PPT"""
        output_path = os.path.join(self.test_dir, "test_md.pptx")
        
        self.skill.create_from_markdown(
            self.sample_markdown,
            output_path,
            template="default",
            title_slide=True
        )
        
        self.assertTrue(os.path.exists(output_path))
        
        # 验证文件可以正常打开
        prs = Presentation(output_path)
        self.assertGreater(len(prs.slides), 0)
    
    def test_create_from_markdown_file(self):
        """测试从Markdown文件创建PPT"""
        md_path = os.path.join(self.test_dir, "input.md")
        output_path = os.path.join(self.test_dir, "test_file.pptx")
        
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write(self.sample_markdown)
        
        self.skill.create_from_markdown(md_path, output_path, template="blue")
        
        self.assertTrue(os.path.exists(output_path))
    
    def test_create_from_json_string(self):
        """测试从JSON字符串创建PPT"""
        output_path = os.path.join(self.test_dir, "test_json.pptx")
        
        self.skill.create_from_json(self.sample_json, output_path, template="green")
        
        self.assertTrue(os.path.exists(output_path))
        
        prs = Presentation(output_path)
        self.assertGreaterEqual(len(prs.slides), 2)
    
    def test_create_from_json_file(self):
        """测试从JSON文件创建PPT"""
        json_path = os.path.join(self.test_dir, "input.json")
        output_path = os.path.join(self.test_dir, "test_json_file.pptx")
        
        with open(json_path, 'w', encoding='utf-8') as f:
            f.write(self.sample_json)
        
        self.skill.create_from_json(json_path, output_path)
        
        self.assertTrue(os.path.exists(output_path))
    
    def test_themes(self):
        """测试所有内置主题"""
        for theme_name in self.skill.THEMES.keys():
            skill = PPTXSkill()
            skill.set_theme(theme_name)
            
            output_path = os.path.join(self.test_dir, f"test_{theme_name}.pptx")
            skill.add_title_slide(f"{theme_name} Theme", "Test")
            skill.save(output_path)
            
            self.assertTrue(os.path.exists(output_path))
    
    def test_save(self):
        """测试保存功能"""
        output_path = os.path.join(self.test_dir, "test_save.pptx")
        
        self.skill.add_slide(title="保存测试", content=["内容"])
        self.skill.save(output_path)
        
        self.assertTrue(os.path.exists(output_path))
        self.assertGreater(os.path.getsize(output_path), 0)
    
    def test_layouts(self):
        """测试不同布局"""
        layouts = ["title", "title_and_content", "title_only", "blank"]
        
        for layout in layouts:
            skill = PPTXSkill()
            slide = skill.add_slide(title=f"Layout: {layout}", layout=layout)
            self.assertIsNotNone(slide)


class TestParsers(unittest.TestCase):
    """解析器测试类"""
    
    def setUp(self):
        if not PPTX_AVAILABLE:
            self.skipTest("python-pptx not installed")
        self.skill = PPTXSkill()
    
    def test_parse_empty_markdown(self):
        """测试空Markdown解析"""
        result = self.skill.parse_markdown("")
        self.assertEqual(len(result), 0)
    
    def test_parse_markdown_with_images(self):
        """测试带图片的Markdown解析"""
        md = """# 图片测试

![描述1](image1.png)
![描述2](image2.jpg)

---

# 另一页

正常内容
"""
        result = self.skill.parse_markdown(md)
        self.assertEqual(len(result), 2)
        self.assertEqual(len(result[0]["images"]), 2)
        self.assertIn("image1.png", result[0]["images"])
    
    def test_parse_markdown_numbered_list(self):
        """测试编号列表解析"""
        md = """# 列表测试

1. 第一项
2. 第二项
3. 第三项
"""
        result = self.skill.parse_markdown(md)
        self.assertEqual(len(result[0]["content"]), 3)
        self.assertIn("第一项", result[0]["content"][0])


class TestIntegration(unittest.TestCase):
    """集成测试类"""
    
    @classmethod
    def setUpClass(cls):
        cls.test_dir = tempfile.mkdtemp()
    
    @classmethod
    def tearDownClass(cls):
        import shutil
        shutil.rmtree(cls.test_dir, ignore_errors=True)
    
    def setUp(self):
        if not PPTX_AVAILABLE:
            self.skipTest("python-pptx not installed")
    
    def test_full_workflow(self):
        """测试完整工作流程"""
        skill = PPTXSkill()
        
        # 添加标题页
        skill.add_title_slide("年度报告", "2026年度总结")
        
        # 添加内容页
        skill.add_slide(
            title="业绩概览",
            content=["收入增长20%", "用户增长35%", "利润增长15%"]
        )
        
        skill.add_slide(
            title="未来规划",
            content=["产品创新", "市场扩展", "团队建设"]
        )
        
        # 保存
        output_path = os.path.join(self.test_dir, "workflow_test.pptx")
        skill.save(output_path)
        
        # 验证
        self.assertTrue(os.path.exists(output_path))
        prs = Presentation(output_path)
        self.assertEqual(len(prs.slides), 3)


def run_tests():
    """运行测试"""
    loader = unittest.TestLoader()
    suite = unittest.TestSuite()
    
    suite.addTests(loader.loadTestsFromTestCase(TestPPTXSkill))
    suite.addTests(loader.loadTestsFromTestCase(TestParsers))
    suite.addTests(loader.loadTestsFromTestCase(TestIntegration))
    
    runner = unittest.TextTestRunner(verbosity=2)
    result = runner.run(suite)
    
    return result.wasSuccessful()


if __name__ == '__main__':
    success = run_tests()
    sys.exit(0 if success else 1)
